"""
pptx_generator.py — python-pptx 기반 PPTX 생성기
FPOF Document Converter

파이프라인:
  themes/{template}.potx 로드
  → 브랜드 컬러 XML 레벨 오버라이드
  → 슬라이드 타입별 레이아웃 생성 (cover / section / content / table / kpi / closing)
  → 저장
"""
from __future__ import annotations
import os
import re
import copy
from lxml import etree
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ..core.content_model import (
    DocumentContent, Section, Block, BlockType, BrandTheme, TableRow
)

# ── 상수 ────────────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

MARGIN_L = Inches(0.8)
MARGIN_R = Inches(0.8)
MARGIN_T = Inches(0.5)
CONTENT_W = W - MARGIN_L - MARGIN_R

# ── 색상 헬퍼 ─────────────────────────────────────────────────────────
def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _hex6(hex_str: str) -> str:
    return hex_str.lstrip('#').upper()


# ── 테마 컬러 XML 오버라이드 ──────────────────────────────────────────
def _override_theme_colors(prs: Presentation, brand: BrandTheme) -> None:
    """
    slide_master의 theme XML에서 dk1/dk2/lt1/lt2/accent1~6 컬러를
    브랜드 컬러로 교체한다.
    """
    color_map = {
        "dk1":     _hex6(brand.secondary),    # black
        "dk2":     _hex6(brand.secondary),
        "lt1":     _hex6(brand.background),   # white
        "lt2":     "F5F5F5",
        "accent1": _hex6(brand.primary),      # vivid orange
        "accent2": _hex6(brand.accent1),      # electric blue
        "accent3": _hex6(brand.accent2),      # hot pink
        "accent4": _hex6(brand.accent3),      # acid green
        "accent5": "444444",
        "accent6": "888888",
        "hlink":   _hex6(brand.accent1),
        "folHlink": "666666",
    }

    for master in prs.slide_masters:
        # theme XML 위치: master.element → r:relationship → theme xml
        theme_elem = master.element.find(
            './/{http://schemas.openxmlformats.org/drawingml/2006/main}theme'
        )
        if theme_elem is None:
            # theme는 별도 part에 있음 — slide_master.theme_color_map 접근
            try:
                theme_elem = master.theme_color_map._element
            except AttributeError:
                pass

        if theme_elem is None:
            continue

        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        for name, hex_val in color_map.items():
            # <a:dk1>, <a:accent1> 등 찾기
            elems = theme_elem.findall(f'.//{{{ns}}}{name}')
            for elem in elems:
                srgb = elem.find(f'{{{ns}}}srgbClr')
                sysClr = elem.find(f'{{{ns}}}sysClr')
                if srgb is not None:
                    srgb.set('val', hex_val)
                elif sysClr is not None:
                    # sysClr → srgbClr 로 교체
                    parent = elem
                    parent.remove(sysClr)
                    new_srgb = etree.SubElement(parent, f'{{{ns}}}srgbClr')
                    new_srgb.set('val', hex_val)


# ── 텍스트 박스 헬퍼 ──────────────────────────────────────────────────
def _add_textbox(slide, left, top, width, height,
                 text, font_size, bold=False,
                 color: str = "#000000",
                 align=PP_ALIGN.LEFT,
                 font_name: str = "Inter") -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    run.font.name = font_name


def _add_rect(slide, left, top, width, height, fill_color: str, line_color: str = None) -> None:
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_color)
    if line_color:
        shape.line.color.rgb = _rgb(line_color)
    else:
        shape.line.fill.background()


# ── 슬라이드 생성 함수들 ─────────────────────────────────────────────

def _slide_cover(prs: Presentation, doc: DocumentContent, brand: BrandTheme) -> None:
    """커버 슬라이드"""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # 배경 전체 채우기
    _add_rect(slide, 0, 0, W, H, brand.secondary)

    # 왼쪽 오렌지 액센트 바
    _add_rect(slide, 0, 0, Inches(0.15), H, brand.primary)

    # 브랜드명 (상단)
    _add_textbox(
        slide, MARGIN_L, Inches(1.0), CONTENT_W, Inches(0.6),
        text=brand.brand_name.upper(),
        font_size=13, bold=False, color=brand.primary,
        align=PP_ALIGN.LEFT,
        font_name=brand.display_font
    )

    # 문서 제목
    _add_textbox(
        slide, MARGIN_L, Inches(1.8), CONTENT_W, Inches(1.8),
        text=doc.title,
        font_size=48, bold=True, color=brand.background,
        align=PP_ALIGN.LEFT,
        font_name=brand.display_font
    )

    # 서브타이틀
    if doc.subtitle:
        _add_textbox(
            slide, MARGIN_L, Inches(3.7), CONTENT_W, Inches(0.6),
            text=doc.subtitle,
            font_size=20, bold=False, color="#CCCCCC",
            align=PP_ALIGN.LEFT,
            font_name=brand.body_font
        )

    # 시즌 + 날짜 (하단)
    meta_parts = []
    if doc.season:
        meta_parts.append(doc.season)
    if doc.date:
        meta_parts.append(str(doc.date))
    if not meta_parts:
        meta_parts.append(datetime.today().strftime("%Y.%m"))
    _add_textbox(
        slide, MARGIN_L, Inches(6.4), CONTENT_W, Inches(0.5),
        text="  ·  ".join(meta_parts),
        font_size=12, bold=False, color="#888888",
        align=PP_ALIGN.LEFT,
        font_name=brand.body_font
    )


def _slide_section(prs: Presentation, section: Section, brand: BrandTheme,
                   section_idx: int) -> None:
    """섹션 구분 슬라이드 (H2)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_rect(slide, 0, 0, W, H, "#111111")
    _add_rect(slide, 0, 0, Inches(0.15), H, brand.primary)

    # 섹션 번호
    _add_textbox(
        slide, MARGIN_L, Inches(2.8), Inches(2), Inches(0.8),
        text=f"{section_idx:02d}",
        font_size=72, bold=True, color=brand.primary,
        align=PP_ALIGN.LEFT,
        font_name=brand.display_font
    )

    # 섹션 제목
    _add_textbox(
        slide, MARGIN_L + Inches(1.8), Inches(3.0), CONTENT_W - Inches(1.8), Inches(1.2),
        text=section.title,
        font_size=36, bold=True, color=brand.background,
        align=PP_ALIGN.LEFT,
        font_name=brand.display_font
    )


def _slide_content(prs: Presentation, title: str, blocks: list[Block],
                   brand: BrandTheme) -> None:
    """일반 콘텐츠 슬라이드 (텍스트 + 불릿)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_rect(slide, 0, 0, W, H, brand.background)
    _add_rect(slide, 0, 0, W, Inches(0.08), brand.primary)

    # 슬라이드 제목
    _add_textbox(
        slide, MARGIN_L, MARGIN_T + Inches(0.1), CONTENT_W, Inches(0.7),
        text=title,
        font_size=24, bold=True, color=brand.secondary,
        align=PP_ALIGN.LEFT,
        font_name=brand.display_font
    )

    # 콘텐츠
    y = Inches(1.3)
    line_h = Inches(0.38)
    max_y  = Inches(6.8)

    for block in blocks:
        if y > max_y:
            break
        if block.type == BlockType.HEADING3:
            _add_textbox(
                slide, MARGIN_L, y, CONTENT_W, line_h,
                text=block.text,
                font_size=16, bold=True, color=brand.primary,
                font_name=brand.display_font
            )
            y += line_h + Inches(0.05)

        elif block.type == BlockType.PARAGRAPH:
            txt = block.text
            if len(txt) > 120:
                txt = txt[:120] + "…"
            _add_textbox(
                slide, MARGIN_L, y, CONTENT_W, line_h,
                text=txt,
                font_size=14, bold=False, color=brand.text,
                font_name=brand.body_font
            )
            y += line_h

        elif block.type == BlockType.BULLET:
            for item in block.items[:8]:
                if y > max_y:
                    break
                _add_textbox(
                    slide, MARGIN_L + Inches(0.2), y, CONTENT_W - Inches(0.2), line_h,
                    text=f"• {item[:100]}",
                    font_size=14, bold=False, color=brand.text,
                    font_name=brand.body_font
                )
                y += line_h

        elif block.type == BlockType.KPI:
            _slide_kpi_inline(slide, y, block, brand)
            y += Inches(0.7)

        else:
            y += Inches(0.1)


def _slide_kpi_inline(slide, y_pos, block: Block, brand: BrandTheme) -> None:
    """콘텐츠 슬라이드 내부 KPI 인라인"""
    _add_rect(slide, MARGIN_L, y_pos, Inches(2.0), Inches(0.6), brand.primary)
    _add_textbox(
        slide, MARGIN_L + Inches(0.1), y_pos, Inches(1.8), Inches(0.6),
        text=block.kpi_value,
        font_size=22, bold=True, color=brand.background,
        align=PP_ALIGN.LEFT,
        font_name=brand.display_font
    )
    if block.kpi_label:
        _add_textbox(
            slide, MARGIN_L + Inches(2.2), y_pos + Inches(0.1),
            CONTENT_W - Inches(2.2), Inches(0.5),
            text=block.kpi_label,
            font_size=14, bold=False, color=brand.text,
            font_name=brand.body_font
        )


def _slide_kpi(prs: Presentation, title: str, kpi_blocks: list[Block],
               brand: BrandTheme) -> None:
    """KPI 전용 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_rect(slide, 0, 0, W, H, brand.background)
    _add_rect(slide, 0, 0, W, Inches(0.08), brand.primary)

    _add_textbox(
        slide, MARGIN_L, MARGIN_T + Inches(0.1), CONTENT_W, Inches(0.7),
        text=title,
        font_size=24, bold=True, color=brand.secondary,
        font_name=brand.display_font
    )

    # KPI 카드 그리드 (최대 6개, 3열 2행)
    card_w = Inches(3.5)
    card_h = Inches(1.8)
    cols   = 3
    gap_x  = Inches(0.3)
    gap_y  = Inches(0.3)
    start_x = MARGIN_L
    start_y = Inches(1.5)

    for idx, block in enumerate(kpi_blocks[:6]):
        col = idx % cols
        row = idx // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        _add_rect(slide, x, y, card_w, card_h, "#F5F5F5")
        _add_rect(slide, x, y, card_w, Inches(0.05), brand.primary)

        _add_textbox(
            slide, x + Inches(0.2), y + Inches(0.15),
            card_w - Inches(0.4), Inches(0.9),
            text=block.kpi_value,
            font_size=36, bold=True, color=brand.primary,
            align=PP_ALIGN.LEFT,
            font_name=brand.display_font
        )
        if block.kpi_label:
            _add_textbox(
                slide, x + Inches(0.2), y + Inches(1.1),
                card_w - Inches(0.4), Inches(0.5),
                text=block.kpi_label,
                font_size=12, bold=False, color=brand.text_light,
                font_name=brand.body_font
            )


def _slide_table(prs: Presentation, title: str, block: Block,
                 brand: BrandTheme) -> None:
    """테이블 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_rect(slide, 0, 0, W, H, brand.background)
    _add_rect(slide, 0, 0, W, Inches(0.08), brand.primary)

    _add_textbox(
        slide, MARGIN_L, MARGIN_T + Inches(0.1), CONTENT_W, Inches(0.7),
        text=title,
        font_size=24, bold=True, color=brand.secondary,
        font_name=brand.display_font
    )

    if not block.rows:
        return

    rows = [r for r in block.rows if r.cells]
    if not rows:
        return

    cols = max(len(r.cells) for r in rows)
    col_w = CONTENT_W / cols
    row_h = Inches(0.45)
    table_top = Inches(1.4)

    for r_idx, row in enumerate(rows[:12]):
        for c_idx, cell in enumerate(row.cells[:cols]):
            x = MARGIN_L + c_idx * col_w
            y = table_top + r_idx * row_h

            if row.is_header:
                _add_rect(slide, x, y, col_w - Inches(0.02), row_h, brand.secondary)
                txt_color = brand.background
                bold = True
            elif r_idx % 2 == 0:
                _add_rect(slide, x, y, col_w - Inches(0.02), row_h, "#F9F9F9")
                txt_color = brand.text
                bold = False
            else:
                txt_color = brand.text
                bold = False

            _add_textbox(
                slide, x + Inches(0.1), y + Inches(0.05),
                col_w - Inches(0.2), row_h - Inches(0.1),
                text=str(cell)[:80],
                font_size=11, bold=bold, color=txt_color,
                font_name=brand.body_font
            )


def _slide_closing(prs: Presentation, doc: DocumentContent, brand: BrandTheme) -> None:
    """클로징 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_rect(slide, 0, 0, W, H, brand.secondary)
    _add_rect(slide, 0, 0, Inches(0.15), H, brand.primary)

    _add_textbox(
        slide, MARGIN_L, Inches(3.0), CONTENT_W, Inches(1.2),
        text=brand.brand_name.upper(),
        font_size=48, bold=True, color=brand.primary,
        align=PP_ALIGN.LEFT,
        font_name=brand.display_font
    )

    _add_textbox(
        slide, MARGIN_L, Inches(4.3), CONTENT_W, Inches(0.6),
        text=doc.title,
        font_size=18, bold=False, color="#AAAAAA",
        align=PP_ALIGN.LEFT,
        font_name=brand.body_font
    )


# ── 메인 생성 함수 ────────────────────────────────────────────────────

def generate(doc: DocumentContent, brand: BrandTheme,
             template: str = "executive",
             output_path: str = None) -> str:
    """
    DocumentContent + BrandTheme → PPTX 파일 생성

    Args:
        doc:         파서 출력
        brand:       브랜드 테마
        template:    executive / creative / report / internal
        output_path: 저장 경로 (None이면 임시 경로 반환)

    Returns:
        생성된 파일 경로
    """
    _here      = os.path.dirname(os.path.abspath(__file__))
    _converter = os.path.dirname(_here)
    theme_path = os.path.join(_converter, "themes", f"{template}.potx")

    if os.path.exists(theme_path):
        prs = Presentation(theme_path)
    else:
        prs = Presentation()

    prs.slide_width  = W
    prs.slide_height = H

    # 테마 컬러 오버라이드
    try:
        _override_theme_colors(prs, brand)
    except Exception:
        pass  # 오버라이드 실패해도 계속 진행

    # ── 슬라이드 생성 ────────────────────────────────────────────
    # 1. 커버
    _slide_cover(prs, doc, brand)

    # 2. 섹션별 슬라이드
    section_counter = 0
    for section in doc.sections:
        if section.title == "(intro)" and not section.blocks:
            continue

        # 섹션 구분 슬라이드 (intro 제외)
        if section.title != "(intro)":
            section_counter += 1
            _slide_section(prs, section, brand, section_counter)

        # 섹션 내 블록을 슬라이드로 분배
        title_for_content = section.title if section.title != "(intro)" else doc.title
        _distribute_blocks(prs, title_for_content, section.blocks, brand)

    # 3. 클로징
    _slide_closing(prs, doc, brand)

    # ── 저장 ─────────────────────────────────────────────────────
    if output_path is None:
        base = os.path.splitext(os.path.basename(doc.source_path))[0]
        output_path = os.path.join(
            os.path.dirname(doc.source_path), "exports",
            f"{base}_{template}.pptx"
        )

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path


def _distribute_blocks(prs: Presentation, section_title: str,
                       blocks: list[Block], brand: BrandTheme) -> None:
    """
    섹션의 블록들을 슬라이드로 나눈다.
    - 테이블 → 테이블 슬라이드
    - KPI 묶음 (≥3) → KPI 슬라이드
    - 나머지 → 콘텐츠 슬라이드 (MAX_ITEMS씩)
    """
    MAX_ITEMS = 8

    # KPI 전용 슬라이드로 분리할 KPI 블록
    kpi_blocks   = [b for b in blocks if b.type == BlockType.KPI]
    table_blocks = [b for b in blocks if b.type == BlockType.TABLE]
    other_blocks = [b for b in blocks if b.type not in (BlockType.TABLE,)]

    # 테이블 슬라이드
    for t_block in table_blocks:
        _slide_table(prs, section_title, t_block, brand)

    # KPI 전용 슬라이드 (3개 이상이면)
    if len(kpi_blocks) >= 3:
        _slide_kpi(prs, section_title, kpi_blocks, brand)
        # other_blocks에서 KPI 제거
        other_blocks = [b for b in other_blocks if b.type != BlockType.KPI]

    # 일반 콘텐츠 — MAX_ITEMS씩 슬라이드로 나눔
    for start in range(0, max(1, len(other_blocks)), MAX_ITEMS):
        chunk = other_blocks[start:start + MAX_ITEMS]
        if not chunk:
            break
        title = section_title
        if start > 0:
            title = f"{section_title} (계속)"
        _slide_content(prs, title, chunk, brand)
